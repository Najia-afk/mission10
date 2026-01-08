import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# --- 1. Configuration & Assets ---
os.makedirs("assets", exist_ok=True)

# Brand Palette (Vibrant Startup)
COLOR_BG = RGBColor(255, 255, 255)       # White
COLOR_TEXT_MAIN = RGBColor(40, 40, 50)   # Dark Slate
COLOR_ACCENT_1 = RGBColor(138, 43, 226)  # Blue Violet
COLOR_ACCENT_2 = RGBColor(255, 20, 147)  # Deep Pink
COLOR_ACCENT_3 = RGBColor(0, 200, 255)   # Cyan

# --- 2. Data & Charts Generation ---

# Backlog Data
backlog_data = [
    {"ID": "US01", "Titre": "Connexion mail", "Priorité": "Must", "Points": 3},
    {"ID": "US02", "Titre": "Capture photo", "Priorité": "Must", "Points": 5},
    {"ID": "US03", "Titre": "Moteur Reco IA", "Priorité": "Must", "Points": 13},
    {"ID": "US04", "Titre": "Virtual Try-on", "Priorité": "Should", "Points": 21},
    {"ID": "US05", "Titre": "Profil Style", "Priorité": "Should", "Points": 5},
    {"ID": "US09", "Titre": "Paiement In-App", "Priorité": "Must", "Points": 13},
]
df_backlog = pd.DataFrame(backlog_data)

# ROI Chart (Light Theme)
mois = np.arange(0, 37)
depenses = 50000 + (12000 * mois) # Simplified model for visual clarity
recettes = np.maximum(0, 25000 * (mois - 6))

fig_roi = go.Figure()
fig_roi.add_trace(go.Scatter(x=mois, y=depenses, name='Investment', line=dict(color='#FF1493', width=4)))
fig_roi.add_trace(go.Scatter(x=mois, y=recettes, name='Revenue', line=dict(color='#8A2BE2', width=4)))
fig_roi.add_vline(x=18, line_dash="dash", line_color="gray", annotation_text="Break-even")

fig_roi.update_layout(
    title='Path to Profitability',
    template='plotly_white',
    font=dict(family="Arial", size=14, color="#333"),
    xaxis=dict(showgrid=False),
    yaxis=dict(showgrid=True, gridcolor='#eee'),
    legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
)
fig_roi.write_image("assets/roi_chart_light.png", scale=2)

# Risk Radar (Light Theme)
risks = [
    {"Cat": "Data Privacy", "Crit": 8}, {"Cat": "Legal/GDPR", "Crit": 7}, 
    {"Cat": "AI Ethics", "Crit": 5}, {"Cat": "Tech Scalability", "Crit": 4}, 
    {"Cat": "Market Adoption", "Crit": 6}, {"Cat": "Talent Retention", "Crit": 5}
]
df_risks = pd.DataFrame(risks)
fig_radar = go.Figure(data=go.Scatterpolar(
    r=df_risks['Crit'], theta=df_risks['Cat'], fill='toself', 
    line=dict(color='#8A2BE2'), fillcolor='rgba(138, 43, 226, 0.2)'
))
fig_radar.update_layout(
    title='Risk Mitigation Profile',
    template='plotly_white',
    polar=dict(radialaxis=dict(visible=True, range=[0, 10], linecolor='#eee')),
    font=dict(family="Arial", size=12)
)
fig_radar.write_image("assets/risk_radar_light.png", scale=2)

# --- 3. Premium PPTX Generation ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_gradient_bar(slide, top=0, height=0.15):
    """Adds a decorative gradient bar at the top."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(top), prs.slide_width, Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT_1
    shape.line.fill.background()

def add_title(slide, text, subtext=None):
    """Adds a premium title and subtitle."""
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_1
    p.font.name = 'Arial Black'
    
    # Subtitle
    if subtext:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(10), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtext
        p2.font.size = Pt(18)
        p2.font.color.rgb = COLOR_TEXT_MAIN
        p2.font.name = 'Arial'

def add_card(slide, left, top, width, height, title, content, icon_char=None):
    """Adds a content card with a shadow effect."""
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 245, 250)
    shape.line.color.rgb = RGBColor(230, 230, 230)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(left+0.2), Inches(top+0.2), Inches(width-0.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_ACCENT_2
    
    # Content
    txBox2 = slide.shapes.add_textbox(Inches(left+0.2), Inches(top+0.8), Inches(width-0.4), Inches(height-1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = content
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT_MAIN

# --- SLIDE 1: COVER ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
# Split screen background
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), 0, Inches(7.33), Inches(7.5))
slide.shapes.add_picture("assets/hero.png", Inches(6), 0, height=Inches(7.5)) # Crop/Fit logic simplified

add_gradient_bar(slide, top=7.35)
# Text
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(5), Inches(2))
p = txBox.text_frame.paragraphs[0]
p.text = "FASHION\nINSTA."
p.font.size = Pt(80)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT_1
p.font.name = 'Arial Black'

txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(5), Inches(1))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Your Personal AI Stylist."
p2.font.size = Pt(24)
p2.font.color.rgb = COLOR_TEXT_MAIN

# --- SLIDE 2: THE PARADOX (Problem) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "The Closet Paradox", "Why fashion is broken today.")

add_card(slide, 0.5, 2, 3.8, 4, "Choice Paralysis", "Users spend 90 mins/week deciding what to wear, yet wear only 20% of their wardrobe.")
add_card(slide, 4.7, 2, 3.8, 4, "The Return Nightmare", "30% of online purchases are returned. Poor fit and style mismatch cost the industry billions.")
add_card(slide, 8.9, 2, 3.8, 4, "Generic Experience", "E-commerce recommendations are based on 'others bought', not 'what fits YOU'.")

# --- SLIDE 3: THE SOLUTION (Mockup) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "Meet Your AI Stylist", "Hyper-personalized fashion at your fingertips.")

# Mockup Center
slide.shapes.add_picture("assets/mockup.png", Inches(4.5), Inches(1.5), height=Inches(5.5))

# Features Left
add_card(slide, 0.5, 2.5, 3.5, 1.5, "Virtual Try-On", "See it on YOU before buying. Powered by Generative AI.")
add_card(slide, 0.5, 4.5, 3.5, 1.5, "Smart Wardrobe", "Digitize your closet. Mix & match instantly.")

# Features Right
add_card(slide, 9.3, 2.5, 3.5, 1.5, "Style DNA", "AI learns your taste, body shape, and vibe.")
add_card(slide, 9.3, 4.5, 3.5, 1.5, "Eco-Score", "Make sustainable choices with real-time impact tracking.")

# --- SLIDE 4: MARKET VISION ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "A $1.5 Trillion Opportunity", "Riding the wave of Fashion Tech.")

# Big Numbers
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(3), Inches(2))
p = txBox.text_frame.paragraphs[0]
p.text = "$1.5T"
p.font.size = Pt(96)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT_3
p.alignment = PP_ALIGN.CENTER
txBox.text_frame.add_paragraph().text = "Global Fashion Market"

txBox2 = slide.shapes.add_textbox(Inches(5), Inches(2.5), Inches(3), Inches(2))
p = txBox2.text_frame.paragraphs[0]
p.text = "25%"
p.font.size = Pt(96)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT_2
p.alignment = PP_ALIGN.CENTER
txBox2.text_frame.add_paragraph().text = "CAGR AI in Fashion"

txBox3 = slide.shapes.add_textbox(Inches(9), Inches(2.5), Inches(3), Inches(2))
p = txBox3.text_frame.paragraphs[0]
p.text = "Gen-Z"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT_1
p.alignment = PP_ALIGN.CENTER
txBox3.text_frame.add_paragraph().text = "Digital Native Target"

# --- SLIDE 5: PRODUCT STRATEGY (MoSCoW) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "Product Roadmap", "Prioritized for maximum impact (MoSCoW).")

# Visual Timeline/Blocks
shape_must = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(4), Inches(4))
shape_must.fill.solid()
shape_must.fill.fore_color.rgb = COLOR_ACCENT_1
shape_must.text_frame.text = "MUST HAVE (MVP)\n\n• AI Reco Engine\n• Photo Capture\n• GDPR Compliance\n• Secure Auth"
shape_must.text_frame.paragraphs[0].font.bold = True
shape_must.text_frame.paragraphs[0].font.size = Pt(20)

shape_should = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.7), Inches(2.5), Inches(4), Inches(4))
shape_should.fill.solid()
shape_should.fill.fore_color.rgb = COLOR_ACCENT_2
shape_should.text_frame.text = "SHOULD HAVE (V1)\n\n• Virtual Try-On\n• Style Profiling\n• Social Sharing"
shape_should.text_frame.paragraphs[0].font.bold = True
shape_should.text_frame.paragraphs[0].font.size = Pt(20)

shape_could = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.9), Inches(2.5), Inches(4), Inches(4))
shape_could.fill.solid()
shape_could.fill.fore_color.rgb = COLOR_ACCENT_3
shape_could.text_frame.text = "COULD HAVE (Scale)\n\n• B2B Retailer API\n• Marketplace\n• Advanced Gamification"
shape_could.text_frame.paragraphs[0].font.bold = True
shape_could.text_frame.paragraphs[0].font.size = Pt(20)

# --- SLIDE 6: FINANCIALS ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "Path to Profitability", "Sustainable growth model.")
slide.shapes.add_picture("assets/roi_chart_light.png", Inches(1), Inches(2), width=Inches(11))

# --- SLIDE 7: RISKS & ETHICS ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "De-risked Execution", "Proactive management of Data & Ethics.")
slide.shapes.add_picture("assets/risk_radar_light.png", Inches(1), Inches(2), width=Inches(6))

# Text details
add_card(slide, 7.5, 2.5, 5, 1.5, "GDPR & Privacy", "Full compliance with CNIL register. AES-256 Encryption. User consent first.")
add_card(slide, 7.5, 4.5, 5, 1.5, "Ethical AI", "Bias monitoring to ensure fair representation across all body types and ethnicities.")

# --- SLIDE 9: AGILE GOVERNANCE (Methodology) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "Agile Governance", "SCRUM Framework & Team Roles.")

# Roles
add_card(slide, 0.5, 2.5, 4, 1.5, "The Squad", "• Product Owner (Alicia)\n• Scrum Master\n• Lead Data Scientist\n• DevOps Engineer")

# Ceremonies
add_card(slide, 4.7, 2.5, 4, 1.5, "Ceremonies", "• Daily Stand-up (15min)\n• Sprint Planning (3 weeks)\n• Review & Retro")

# Tools
add_card(slide, 8.9, 2.5, 4, 1.5, "Tools", "• Jira (Backlog)\n• GitHub (Version Control)\n• Azure DevOps (CI/CD)\n• MLflow (Tracking)")

# --- SLIDE 10: DETAILED BUDGET (Resources) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "Resource Allocation", "CAPEX vs OPEX Breakdown.")

# Budget Table Visual (Simplified)
shape = slide.shapes.add_table(3, 3, Inches(1), Inches(2.5), Inches(11.33), Inches(3))
table = shape.table
table.columns[0].width = Inches(4)
table.columns[1].width = Inches(3)
table.columns[2].width = Inches(4)

# Headers
table.cell(0, 0).text = "Category"
table.cell(0, 1).text = "Cost Type"
table.cell(0, 2).text = "Estimated Amount"

# Rows
table.cell(1, 0).text = "Development (Man-Days)"
table.cell(1, 1).text = "CAPEX (One-off)"
table.cell(1, 2).text = "€125,000 (MVP)"

table.cell(2, 0).text = "Cloud Infrastructure (Azure)"
table.cell(2, 1).text = "OPEX (Monthly)"
table.cell(2, 2).text = "€1,500 / month"

# Style Table
for i in range(3):
    for j in range(3):
        cell = table.cell(i, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(245, 245, 250)
        cell.text_frame.paragraphs[0].font.size = Pt(18)
        cell.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN

# --- SLIDE 11: DATA PRIVACY & ETHICS (Compliance) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "Compliance & Ethics", "Trust is our currency.")

# CNIL Register
add_card(slide, 0.5, 2.5, 5.5, 2, "CNIL Register (RGPD)", "• Purpose: AI Personalization\n• Data: Photos, Size, Purchase History\n• Retention: 3 years active\n• Rights: Access, Rectification, Erasure")

# Ethics & Bias
add_card(slide, 6.5, 2.5, 5.5, 2, "Ethical AI Framework", "• Fairness: Balanced datasets (Ethnicity/Body Type)\n• Transparency: Explainable AI (XAI)\n• Green AI: Optimized inference for low carbon footprint")

# --- SLIDE 12: VISUAL ROADMAP (Timeline) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "Execution Timeline", "Key milestones to market domination.")

# Timeline Line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(11.33), Inches(0.1))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_TEXT_MAIN

# Milestones
milestones = [
    {"Label": "Q1: FOUNDATION", "Date": "Sprint 1-2", "Pos": 1.5, "Color": COLOR_TEXT_MAIN},
    {"Label": "MVP READY 🚀", "Date": "Sprint 3 (Month 2)", "Pos": 4.5, "Color": COLOR_ACCENT_2},
    {"Label": "VIRTUAL TRY-ON", "Date": "Sprint 4 (Month 3)", "Pos": 7.5, "Color": COLOR_ACCENT_1},
    {"Label": "V1 RELEASE 🏁", "Date": "Month 4", "Pos": 10.5, "Color": COLOR_ACCENT_3},
]

for m in milestones:
    # Dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(m["Pos"]), Inches(3.8), Inches(0.5), Inches(0.5))
    dot.fill.solid()
    dot.fill.fore_color.rgb = m["Color"]
    dot.line.fill.background()
    
    # Label
    txBox = slide.shapes.add_textbox(Inches(m["Pos"]-1), Inches(2.8), Inches(2.5), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = m["Label"]
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = m["Color"]
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    txBox2 = slide.shapes.add_textbox(Inches(m["Pos"]-1), Inches(4.4), Inches(2.5), Inches(1))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = m["Date"]
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.font.alignment = PP_ALIGN.CENTER

# --- SLIDE 13: CLOUD ARCHITECTURE (Sexy Diagram) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar(slide)
add_title(slide, "Scalable Cloud Architecture", "Powered by Microsoft Azure.")

# Background Visual
slide.shapes.add_picture("assets/architecture.png", Inches(4), Inches(2), height=Inches(5))

# Tech Stack Cards (Overlay)
add_card(slide, 0.5, 2.5, 3, 1.2, "Front-End", "React Native\n(iOS/Android)")
add_card(slide, 0.5, 4.5, 3, 1.2, "API Gateway", "Azure App Service\n(Python/FastAPI)")
add_card(slide, 9.8, 2.5, 3, 1.2, "AI Engine", "Azure Cognitive Services\n+ Custom PyTorch Models")
add_card(slide, 9.8, 4.5, 3, 1.2, "Data Lake", "Azure Blob Storage\n(Images & Metadata)")

# --- SLIDE 14: THE ASK ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
# Full gradient background
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_ACCENT_1

txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9.33), Inches(2))
p = txBox.text_frame.paragraphs[0]
p.text = "Join the Revolution."
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9.33), Inches(1))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Seeking $2M Seed Investment"
p2.font.size = Pt(32)
p2.font.color.rgb = RGBColor(255, 255, 255)
p2.alignment = PP_ALIGN.CENTER

prs.save("Pitch_Deck_Fashion_Insta_Startup.pptx")
print("Startup Pitch Deck generated successfully.")
